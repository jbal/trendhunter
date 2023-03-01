"""Output formatters."""

from typing import List

from pptx import Presentation

from trendhunter.tuples import Article, Context
from trendhunter.utils import (
    format_console,
    resize_image,
    resolve_path,
)


def console(articles: List[Article], context: Context):
    for article in articles:
        format_console(__name__).debug(
            f"""Article ({article.url.rsplit("/")[-1]}) {{
    url: {article.url}
    title: {article.text.title}
    description: {article.text.description[:40]}... (truncated)
    metadata: {{
        id: {article.text.metadata.eid}
        category: {article.text.metadata.cid}
    }}
    image: {{
        url: {article.image.url}
        image: {article.image.image[:40]}... (truncated)
    }}
}}
"""
        )


def slideshow(articles: List[Article], context: Context) -> None:
    path = resolve_path(context.path)

    if path.exists():
        show = Presentation(path)
    else:
        show = Presentation()

    for article in articles:
        slide = show.slides.add_slide(show.slide_layouts[8])
        slide.placeholders[0].text = article.text.title

        # only insert picture if it can be resized properly
        try:
            im = resize_image(article.image, context.size)
        except (IOError, OSError) as e:
            format_console(__name__).error(
                f'Error resizing image from "{article.image.url}". Error: {e}.'
            )
        else:
            slide.placeholders[1].insert_picture(im)

        slide.placeholders[2].text = article.text.description

    show.save(path)
    format_console(__name__).info(
        f'Formatted {len(articles)} articles to "{path}".'
    )
